# -*- coding: utf-8 -*-
"""数据地图 PPT 生成器（python-pptx，全原生可编辑对象）。

输入：省份评分 JSON（schema 见 references/data-schema.md）
输出：可编辑 PPTX —— 地图=自由多边形 / 表格=原生表格 / 文字=文本框
自动生成页：每评分维度一张地图页 + 综合地图页 + 综合排序条形页 + 维度矩阵表页
            +（数据含 依据_* 时）依据与短板页 +（含 _config.rule 时）评级规则页

用法：
    python gen_pptx.py path/to/data.json [out.pptx]
示例：
    python gen_pptx.py ../examples/example_data.json ../examples/output/example.pptx

与 geo_map.py 同目录；省界数据默认 data/china_bound.json（相对本脚本所在目录）。
"""
import json
import io
import math
import os
import sys

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.dml import MSO_LINE
from pptx.oxml.ns import qn

from geo_map import build_map

_D = os.path.dirname(os.path.abspath(__file__))
DEFAULT_BOUND = os.path.join(_D, '..', 'data', 'china_bound.json')

DEFAULT_GRADE_FILL = {'第一档': 'C6EFCE', '第二档': 'E2EFDA', '第三档': 'FFEB9C',
                      '第四档': 'FCE4D6', '第五档': 'FFC7CE'}
GRADE_NUM = {k: str(i + 1) for i, k in enumerate(DEFAULT_GRADE_FILL)}


def num(x):
    x = x or 0
    return int(x) if x == int(x) else x


def set_ea(run, name='微软雅黑'):
    rPr = run._r.get_or_add_rPr()
    ea = rPr.find(qn('a:ea'))
    if ea is None:
        ea = rPr.makeelement(qn('a:ea'), {})
        rPr.append(ea)
    ea.set('typeface', name)


def txt(sl, x, y, w, h, runs, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, wrap=True):
    tb = sl.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = wrap
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    first = True
    for line_runs, a in runs:
        para = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        para.alignment = a if a else align
        for text, size, color, bold in line_runs:
            r = para.add_run()
            r.text = text
            r.font.size = Pt(size)
            r.font.bold = bold
            r.font.color.rgb = RGBColor.from_string(color)
            r.font.name = '微软雅黑'
            set_ea(r)
    return tb


def title_bar(sl, text, sub=''):
    bar = sl.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(0.28), Inches(12.33), Inches(0.62))
    bar.fill.solid()
    bar.fill.fore_color.rgb = RGBColor.from_string('1F4E79')
    bar.line.fill.background()
    tf = bar.text_frame
    tf.margin_left = Inches(0.15)
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = text
    r.font.size = Pt(20)
    r.font.bold = True
    r.font.color.rgb = RGBColor.from_string('FFFFFF')
    r.font.name = '微软雅黑'
    set_ea(r)
    if sub:
        txt(sl, 0.55, 0.98, 12.2, 0.3, [([(sub, 10.5, '666666', False)], PP_ALIGN.LEFT)])


def make_rank(scores):
    order = sorted(scores, key=lambda p: -scores[p])
    ranks, prev, pr = {}, None, 0
    for i, p in enumerate(order, 1):
        if scores[p] != prev:
            pr, prev = i, scores[p]
        ranks[p] = pr
    return order, ranks


def mk_verdicts(scores, note=''):
    od = sorted(scores, key=lambda p: -scores[p])
    tops = '、'.join('%s%s' % (p, num(scores[p])) for p in od[:3])
    bots = '、'.join('%s%s' % (p, num(scores[p])) for p in od[-3:][::-1])
    lines = [([('前三：%s' % tops, 10.5, '333333', False)], PP_ALIGN.LEFT),
             ([('末三：%s' % bots, 10.5, '333333', False)], PP_ALIGN.LEFT)]
    if note:
        lines.append([(note, 9, '888888', False)])
    return lines


# ---------- 地图页（可编辑自由多边形） ----------
def add_map_slide(prs, title, sub, MAP, verdicts, leg_title):
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    title_bar(sl, title, sub)
    ASPECT = MAP['W'] / MAP['H']
    MX, AVAIL_W, AVAIL_H = 0.55, 8.05, 5.85
    MW = min(AVAIL_W, AVAIL_H * ASPECT)
    MH = MW / ASPECT
    MY = 1.45 + (AVAIL_H - MH) / 2
    scale = int(MW * 914400 / MAP['W'])
    px, py, psc = MAP['_proj']
    ipx, ipy, ipsc = MAP['_iproj']
    IW = 0.9
    LX = max(8.75, MX + MW + 0.35)
    scale_i = int(IW * 914400 / MAP['iW'])
    IX, IY = LX, 5.62

    for nm, ring, tgt, fill in MAP['_rings_main']:
        pts = [((x - px) * psc, (y - py) * psc) for x, y in ring]
        ff = sl.shapes.build_freeform(pts[0][0], pts[0][1], scale)
        ff.add_line_segments(pts[1:], close=True)
        shp = ff.convert_to_shape(Emu(int(MX * 914400)), Emu(int(MY * 914400)))
        shp.fill.solid()
        shp.fill.fore_color.rgb = RGBColor.from_string(fill.lstrip('#').upper())
        shp.line.color.rgb = RGBColor.from_string('FFFFFF')
        shp.line.width = Pt(0.75)
    box = sl.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(IX - 0.07), Inches(IY - 0.07),
                              Inches(IW + 0.14), Inches(IW * MAP['iH'] / MAP['iW'] + 0.14))
    box.fill.solid()
    box.fill.fore_color.rgb = RGBColor.from_string('FFFFFF')
    box.line.color.rgb = RGBColor.from_string('808080')
    box.line.width = Pt(0.75)
    txt(sl, IX - 0.07, IY + IW * MAP['iH'] / MAP['iW'] + 0.09, IW + 0.14, 0.2,
        [([('南海诸岛', 8, '666666', False)], PP_ALIGN.CENTER)])
    for ring, jd in MAP['_rings_inset']:
        pts = [((x - ipx) * ipsc, (y - ipy) * ipsc) for x, y in ring]
        ff = sl.shapes.build_freeform(pts[0][0], pts[0][1], scale_i)
        ff.add_line_segments(pts[1:], close=True)
        shp = ff.convert_to_shape(Emu(int(IX * 914400)), Emu(int(IY * 914400)))
        if jd:
            shp.fill.background()
            shp.line.color.rgb = RGBColor.from_string('8C8C8C')
            shp.line.width = Pt(0.75)
            shp.line.dash_style = MSO_LINE.DASH
        else:
            shp.fill.solid()
            shp.fill.fore_color.rgb = RGBColor.from_string('EFEFEF')
            shp.line.color.rgb = RGBColor.from_string('FFFFFF')
            shp.line.width = Pt(0.5)

    for L in MAP['labels']:
        cx = MX + L['x'] * MW / MAP['W']
        cy = MY + L['y'] * MW / MAP['W']
        if L['t'] is False:
            tb = sl.shapes.add_textbox(Inches(cx - 0.4), Inches(cy - 0.09), Inches(0.8), Inches(0.2))
            tf = tb.text_frame
            tf.word_wrap = False
            tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
            p1 = tf.paragraphs[0]
            p1.alignment = PP_ALIGN.CENTER
            r = p1.add_run()
            r.text = L['n']
            r.font.size = Pt(7)
            r.font.color.rgb = RGBColor.from_string('999999')
            r.font.name = '微软雅黑'
            set_ea(r)
            continue
        dark = L['b'] <= 1
        col = 'FFFFFF' if dark else '17406B'
        fs1 = 11 if L['fs'] >= 12 else 10
        tb = sl.shapes.add_textbox(Inches(cx - 0.55), Inches(cy - 0.24), Inches(1.1), Inches(0.48))
        tf = tb.text_frame
        tf.word_wrap = False
        tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
        p1 = tf.paragraphs[0]
        p1.alignment = PP_ALIGN.CENTER
        r = p1.add_run()
        r.text = L['n']
        r.font.size = Pt(fs1)
        r.font.bold = True
        r.font.color.rgb = RGBColor.from_string(col)
        r.font.name = '微软雅黑'
        set_ea(r)
        p2 = tf.add_paragraph()
        p2.alignment = PP_ALIGN.CENTER
        r = p2.add_run()
        r.text = str(L['s'])
        r.font.size = Pt(fs1 - 2)
        r.font.color.rgb = RGBColor.from_string(col)
        r.font.name = '微软雅黑'
        set_ea(r)

    txt(sl, LX, 1.5, 4.1, 0.3, [([('结论速览', 12, '1F4E79', True)], PP_ALIGN.LEFT)])
    txt(sl, LX, 1.84, 4.15, 1.0, verdicts)
    txt(sl, LX, 3.1, 4.1, 0.3, [([(leg_title, 12, '1F4E79', True)], PP_ALIGN.LEFT)])
    ly = 3.46
    for item in MAP['leg']:
        sq = sl.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(LX), Inches(ly), Inches(0.24), Inches(0.24))
        sq.fill.solid()
        sq.fill.fore_color.rgb = RGBColor.from_string(item['c'].lstrip('#').upper())
        sq.line.color.rgb = RGBColor.from_string('B0B0B0')
        sq.line.width = Pt(0.5)
        txt(sl, LX + 0.34, ly + 0.01, 3.6, 0.24, [([(item['t'], 10.5, '333333', False)], PP_ALIGN.LEFT)])
        ly += 0.33
    sq = sl.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(LX), Inches(ly), Inches(0.24), Inches(0.24))
    sq.fill.solid()
    sq.fill.fore_color.rgb = RGBColor.from_string('EFEFEF')
    sq.line.color.rgb = RGBColor.from_string('B0B0B0')
    sq.line.width = Pt(0.5)
    txt(sl, LX + 0.34, ly + 0.01, 3.6, 0.24, [([('非调研省份（灰）', 10.5, '333333', False)], PP_ALIGN.LEFT)])
    return sl


# ---------- 综合排序条形页 ----------
def add_bar_slide(prs, P, title, subtitle, groups, totals):
    """totals = 各省原始分总和（各 group 分相加）；条形分段=各 group 原始分，刻度与数字均按 totals 自洽。"""
    order, ranks = make_rank(totals)
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    title_bar(sl, title, subtitle)
    # 分色：前 group 深蓝，后续 group 浅蓝递进（最多 6 段足够）
    COLORS = ['1F4E79', '5B9BD5', '9DC3E6', '2E75B6', '8FAADC', 'BDD7EE']
    if len(groups) > len(COLORS):
        raise ValueError('条形页最多支持 %d 个分色维度' % len(COLORS))
    full = max((totals[p] for p in order), default=1) or 1
    BX, BY, BWD = 2.3, 1.55, 9.0
    RH = 0.355
    for v in range(0, 21, 5):
        gx = BX + BWD * v / 20
        ln = sl.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(gx), Inches(BY - 0.08), Inches(0.008), Inches(len(order) * RH + 0.1))
        ln.fill.solid()
        ln.fill.fore_color.rgb = RGBColor.from_string('E3E9F0')
        ln.line.fill.background()
        txt(sl, gx - 0.2, BY + len(order) * RH + 0.06, 0.4, 0.22, [([(str(int(full * v / 20)) if v < 20 else str(int(full)), 8.5, '999999', False)], PP_ALIGN.CENTER)])
    keys = [g['keys'] for g in groups]
    for i, p in enumerate(order):
        y = BY + i * RH
        txt(sl, 0.5, y + 0.03, 1.7, 0.28, [([('%d. %s' % (ranks[p], p), 11, '333333', True)], PP_ALIGN.RIGHT)])
        x = BX
        for gi, gk in enumerate(keys):
            v = sum(P[p].get(k, 0) or 0 for k in gk)
            w = BWD * v / full
            if w > 0.01:
                b1 = sl.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(RH - 0.07))
                b1.fill.solid()
                b1.fill.fore_color.rgb = RGBColor.from_string(COLORS[gi])
                b1.line.fill.background()
            x += w
        txt(sl, x + 0.08, y + 0.02, 0.8, 0.28, [([(str(num(totals[p])), 11, '1F4E79', True)], PP_ALIGN.LEFT)])
    return sl


# ---------- 维度矩阵表页 ----------
def add_matrix_slide(prs, data, order, ranks, totals, groups, grade_fill):
    P = data['省份']
    ratings = sorted({k[3:] for d in P.values() for k in d if k.startswith('评级_')})
    headers = ['省份', '名次', '总分'] + [r.replace('_', '·') for r in ratings] + [g['name'] for g in groups]
    widths = [1.0, 0.7, 0.8] + [1.0] * len(ratings) + [1.1] * len(groups)
    ncol = len(headers)
    tbl = prs.slides.add_slide(prs.slide_layouts[6])
    title_bar(tbl, '维度档位矩阵', '省份 × 各评分维度分项与评级（评级为资料评价，不计分）')
    table = tbl.shapes.add_table(len(order) + 1, ncol, Inches(0.55), Inches(1.45),
                                 Inches(sum(widths)), Inches(0.4 + len(order) * 0.36)).table
    for ci, wd in enumerate(widths):
        table.columns[ci].width = Inches(wd)
    for ci, h in enumerate(headers):
        c = table.cell(0, ci)
        c.fill.solid()
        c.fill.fore_color.rgb = RGBColor.from_string('1F4E79')
        c.vertical_anchor = MSO_ANCHOR.MIDDLE
        tf = c.text_frame
        tf.word_wrap = True
        para = tf.paragraphs[0]
        para.alignment = PP_ALIGN.CENTER
        r = para.add_run()
        r.text = h
        r.font.size = Pt(10)
        r.font.bold = True
        r.font.color.rgb = RGBColor.from_string('FFFFFF')
        r.font.name = '微软雅黑'
        set_ea(r)
    table.rows[0].height = Inches(0.42)
    for ri, p in enumerate(order, 1):
        d = P[p]
        table.rows[ri].height = Inches(0.36)
        vals = [p, str(ranks[p]), str(num(totals[p]))]
        for rk in ratings:
            vals.append(GRADE_NUM.get(d.get('评级_' + rk, ''), '?'))
        for g in groups:
            vals.append(str(num(sum(d.get(k, 0) or 0 for k in g['keys']))))
        for ci, v in enumerate(vals):
            c = table.cell(ri, ci)
            c.vertical_anchor = MSO_ANCHOR.MIDDLE
            tf = c.text_frame
            para = tf.paragraphs[0]
            para.alignment = PP_ALIGN.CENTER
            r = para.add_run()
            r.text = v
            r.font.size = Pt(10.5)
            r.font.name = '微软雅黑'
            set_ea(r)
            r.font.color.rgb = RGBColor.from_string('333333')
            if ci == 0:
                r.font.bold = True
            if 3 <= ci < 3 + len(ratings) and v in ('1', '2', '3', '4', '5'):
                c.fill.solid()
                c.fill.fore_color.rgb = RGBColor.from_string(grade_fill[int(v) - 1])
                r.font.bold = True
                r.font.size = Pt(11)
    return tbl


# ---------- 依据与短板页 ----------
def est_h(text, chars_per_line):
    lines = max(1, math.ceil(len(text) / chars_per_line))
    return min(0.125 * lines + 0.07, 1.1)


def reason_cell(d, g):
    """group 的依据文本 = 各 score_ 键对应的 依据_<维度> 字段拼接（如 score_用户侧 → 依据_用户侧）。"""
    parts = []
    for k in g['keys']:
        v = d.get('依据_' + k[6:], '')
        if v:
            parts.append(''.join(v) if isinstance(v, list) else v)
    return '；'.join(parts)


def add_reason_slides(prs, data, order, groups, per_page=5):
    P = data['省份']
    for si, provs in enumerate([order[i:i + per_page] for i in range(0, len(order), per_page)]):
        sl = prs.slides.add_slide(prs.slide_layouts[6])
        title_bar(sl, '各省评分依据与短板（%d/%d）' % (si + 1, math.ceil(len(order) / per_page)),
                  '依据=分项得分/满分结构 · 短板为扣分主因')
        yw = [0.7] + [2.5] * len(groups) + [1.5]
        y = 1.5
        hrow = 0.42
        hh = sl.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.55), Inches(y), Inches(sum(yw)), Inches(hrow))
        hh.fill.solid()
        hh.fill.fore_color.rgb = RGBColor.from_string('1F4E79')
        hh.line.fill.background()
        heads = ['省份'] + ['%s\n评分依据' % g['name'] for g in groups] + ['关键短板']
        for ci, htxt in enumerate(heads):
            txt(sl, 0.55 + sum(yw[:ci]) + 0.06, y + 0.05, yw[ci] - 0.12, hrow - 0.08,
                [([(htxt, 9, 'FFFFFF', True)], PP_ALIGN.CENTER)])
        y += hrow
        for p in provs:
            d = P[p]
            cells = [p] + [reason_cell(d, g) for g in groups] + [d.get('短板', '')]
            cpl = [8] + [23] * len(groups) + [13]
            rh = max(est_h(str(c), cpl[i]) for i, c in enumerate(cells))
            for ci, ctext in enumerate(cells):
                box = sl.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.55 + sum(yw[:ci])), Inches(y),
                                          Inches(yw[ci]), Inches(rh))
                box.fill.solid()
                box.fill.fore_color.rgb = RGBColor.from_string('FFF8E7' if ci == len(cells) - 1 else ('EAF1F8' if ci == 0 else 'FFFFFF'))
                box.line.color.rgb = RGBColor.from_string('D9D9D9')
                box.line.width = Pt(0.5)
                tf = box.text_frame
                tf.word_wrap = True
                tf.margin_left = tf.margin_right = Inches(0.05)
                tf.margin_top = tf.margin_bottom = Inches(0.02)
                para = tf.paragraphs[0]
                para.alignment = PP_ALIGN.CENTER if ci == 0 else PP_ALIGN.LEFT
                r = para.add_run()
                r.text = str(ctext)
                r.font.size = Pt(7.5 if ci else 9.5)
                r.font.bold = ci == 0
                r.font.color.rgb = RGBColor.from_string('333333' if ci != len(cells) - 1 else '7F6000')
                r.font.name = '微软雅黑'
                set_ea(r)
            y += rh


# ---------- 评级规则页 ----------
def add_rule_slide(prs, rule):
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    title_bar(sl, '评级规则与说明')
    tb = sl.shapes.add_textbox(Inches(0.55), Inches(1.4), Inches(12.2), Inches(5.6))
    tf = tb.text_frame
    tf.word_wrap = True
    for li, line in enumerate(str(rule).split('\n')):
        para = tf.paragraphs[0] if li == 0 else tf.add_paragraph()
        r = para.add_run()
        r.text = line
        r.font.size = Pt(11)
        r.font.color.rgb = RGBColor.from_string('444444')
        r.font.name = '微软雅黑'
        set_ea(r)
        para.space_after = Pt(3)
    return sl


def build(data, out='out.pptx'):
    P = data['省份']
    cfg = data.get('_config', {})
    if not P:
        raise ValueError('数据 JSON 缺少"省份"节')
    groups = cfg.get('groups')
    if not groups:
        # 自动发现：所有 score_ 前缀字段各作一组
        keys = [k for k in next(iter(P.values())) if k.startswith('score_')]
        groups = {k[6:]: {'keys': [k]} for k in keys}
    groups = [{'name': name, 'keys': v['keys'], 'note': v.get('note', ''), 'cap': v.get('cap', 10)}
              for name, v in groups.items()]

    def group_score(p, g):
        return sum(P[p].get(k, 0) or 0 for k in g['keys'])

    wcfg = cfg.get('综合', {})
    weights = wcfg.get('weights')
    if weights:
        combined = {p: round(sum(weights[g['name']] * (group_score(p, g) / (g['cap'] or 1)) for g in groups), 2) for p in P}
    else:
        combined = {p: round(sum(group_score(p, g) for g in groups), 2) for p in P}

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    title = cfg.get('title', '数据评分图谱')
    subtitle = cfg.get('subtitle', '')
    note = cfg.get('note', '')
    grade_fill = [cfg.get('等级_填充', {}).get(k, v) for k, v in DEFAULT_GRADE_FILL.items()]

    for g in groups:
        sc = {p: group_score(p, g) for p in P}
        od, rk = make_rank(sc)
        MAP = build_map(set(P), sc, od, rk, bound=cfg.get('bound', DEFAULT_BOUND))
        add_map_slide(prs, '%s · %s' % (title, g['name']), subtitle + (' ｜ ' + g['note'] if g['note'] else ''),
                      MAP, mk_verdicts(sc, note), '图例（按%s分）' % g['name'])

    od_c, rk_c = make_rank(combined)
    MAP_C = build_map(set(P), combined, od_c, rk_c, bound=cfg.get('bound', DEFAULT_BOUND))
    add_map_slide(prs, '%s · 综合' % title,
                  subtitle + (' ｜ ' + wcfg.get('note', '') if wcfg.get('note') else ''),
                  MAP_C, mk_verdicts(combined, note), '图例（按综合分）')

    raw_total = {p: round(sum(group_score(p, g) for g in groups), 2) for p in P}
    add_bar_slide(prs, P, '综合排序（结论先行）', '横条按原始总分降序（各维度分叠加），并列名次相同',
                  groups, raw_total)
    add_matrix_slide(prs, data, od_c, rk_c, combined, groups, grade_fill)
    if any('短板' in d for d in P.values()) or any(reason_cell(d, g) for d in P.values() for g in groups):
        add_reason_slides(prs, data, od_c, groups)
    if cfg.get('rule'):
        add_rule_slide(prs, cfg['rule'])

    tmp = out + '.tmp'
    prs.save(tmp)
    os.replace(tmp, out)
    print('saved %s | slides: %d' % (out, len(prs.slides._sldIdLst)))


if __name__ == '__main__':
    if len(sys.argv) < 2:
        print(__doc__)
        sys.exit(1)
    data = json.load(io.open(sys.argv[1], encoding='utf-8'))
    out = sys.argv[2] if len(sys.argv) > 2 else (sys.argv[1][:-5] + '.pptx')
    build(data, out)
